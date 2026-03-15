"""Generate a dark-themed conference presentation on prompt injection.

Composes structured research data from src.research modules into a
~15-slide .pptx with dark backgrounds, speaker notes, embedded figures,
slide numbers, visual variety, and accent-colored labels.

Usage:
    python -m src.presentation.generate
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from src.research.taxonomy import ATTACK_TAXONOMY
from src.research.risks import RISK_CATEGORIES
from src.research.detection import DETECTION_TECHNIQUES
from src.research.prevention import PREVENTION_STRATEGIES
from src.research.references import REFERENCES

PROJECT_ROOT = Path(__file__).resolve().parent.parent.parent
OUTPUT_DIR = PROJECT_ROOT / "output"
FIGURES_DIR = PROJECT_ROOT / "results" / "figures"

# Theme colors
BG_COLOR = RGBColor(30, 30, 30)
TEXT_COLOR = RGBColor(230, 230, 230)
HEADING_COLOR = RGBColor(100, 180, 255)
ACCENT_COLOR = RGBColor(255, 170, 60)
SUBTITLE_COLOR = RGBColor(180, 180, 180)
CRITICAL_COLOR = RGBColor(255, 80, 80)
HIGH_COLOR = RGBColor(255, 170, 60)
DIM_COLOR = RGBColor(140, 140, 140)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

AUTHOR = "Edward Griggs"


def _set_dark_bg(slide):
    """Apply solid dark background fill to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def _add_textbox(slide, left, top, width, height):
    """Add a textbox and return it."""
    return slide.shapes.add_textbox(left, top, width, height)


def _set_text(textbox, text, *, font_size=18, color=TEXT_COLOR, bold=False, alignment=PP_ALIGN.LEFT):
    """Set text on a textbox's first paragraph."""
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment


def _add_heading(slide, text, top=Inches(0.4)):
    """Add a heading textbox at the top of a slide."""
    tb = _add_textbox(slide, Inches(0.8), top, Inches(11.5), Inches(0.8))
    _set_text(tb, text, font_size=32, color=HEADING_COLOR, bold=True)
    return tb


def _add_bullets(slide, items, *, left=Inches(0.8), top=Inches(1.5), width=Inches(11.5),
                 height=Inches(5.0), font_size=18, color=TEXT_COLOR, max_items=None):
    """Add a bulleted list to a slide."""
    tb = _add_textbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    display_items = items[:max_items] if max_items else items
    for i, item in enumerate(display_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022 {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def _add_rich_bullets(slide, items, *, left=Inches(0.8), top=Inches(1.5), width=Inches(11.5),
                      height=Inches(5.0), font_size=17, color=TEXT_COLOR,
                      label_color=ACCENT_COLOR):
    """Add bullets with accent-colored labels.

    items: list of (label, description) tuples.
    """
    tb = _add_textbox(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (label, desc) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run_bullet = _make_run(p, "\u2022 ")
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = color
        run_label = _make_run(p, label)
        run_label.font.size = Pt(font_size)
        run_label.font.color.rgb = label_color
        run_label.font.bold = True
        run_desc = _make_run(p, f" \u2014 {desc}")
        run_desc.font.size = Pt(font_size)
        run_desc.font.color.rgb = color
        p.space_after = Pt(8)


def _add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def _make_run(paragraph, text):
    """Create a run with text on a paragraph (python-pptx add_run takes no args)."""
    run = paragraph.add_run()
    run.text = text
    return run


def _add_slide_number(slide, number):
    """Add a slide number in the bottom-right corner."""
    tb = _add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4))
    _set_text(tb, str(number), font_size=10, color=DIM_COLOR, alignment=PP_ALIGN.RIGHT)


def _add_figure(slide, filename, left=Inches(1.0), top=Inches(1.5),
                width=Inches(5.5), height=None):
    """Embed a figure image on a slide. Returns True if successful."""
    path = FIGURES_DIR / filename
    if not path.exists():
        print(f"  WARNING: Figure not found: {path}")
        return False
    if height:
        slide.shapes.add_picture(str(path), left, top, width, height)
    else:
        slide.shapes.add_picture(str(path), left, top, width)
    print(f"  Embedded figure: {filename}")
    return True


def generate_presentation() -> Path:
    """Build and save the presentation, returning the output path."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # blank layout
    slide_num = 0

    # --- Slide 1: Title ---
    print("Creating slide 1: Title")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    tb = _add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(2))
    _set_text(tb, "Prompt Injection Attacks in Large Language Models",
              font_size=36, color=HEADING_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    tb2 = _add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1))
    _set_text(tb2, "A Survey of Attack Taxonomies, Risks, Detection, and Prevention",
              font_size=22, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
    # Author name
    tb3 = _add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.6))
    _set_text(tb3, AUTHOR, font_size=20, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)
    # Affiliation
    tb4 = _add_textbox(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.6))
    _set_text(tb4, "Christopher Newport University \u2022 COVACCI Cybersecurity Undergraduate Research",
              font_size=16, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
    _add_slide_number(slide, slide_num)

    # --- Slide 2: Agenda / Overview ---
    print("Creating slide 2: Agenda")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Presentation Overview")
    _add_bullets(slide, [
        "What is prompt injection?",
        "Attack taxonomy \u2014 direct and indirect vectors",
        "Real-world risks and impact categories",
        "Detection techniques and tools",
        "Prevention strategies and defense architecture",
        "Key references and future directions",
    ])
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~30s] Welcome audience. Outline the four research pillars: taxonomy, risks, detection, prevention. Set expectation for 5-8 min talk.")

    # --- Slide 3: What is Prompt Injection? (with stat callout) ---
    print("Creating slide 3: Introduction with stat callout")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "What is Prompt Injection?")
    # Bullets on the left (narrower)
    _add_bullets(slide, [
        "Adversarial inputs that hijack LLM behavior by overriding system instructions",
        "Analogous to SQL injection \u2014 exploits the blurred boundary between code and data",
        "Affects all major LLM platforms: ChatGPT, Claude, Gemini, Copilot",
        "Two primary vectors: direct injection (user input) and indirect injection (external data)",
    ], left=Inches(0.8), width=Inches(8.0), font_size=17)
    # Stat callout on the right
    tb_num = _add_textbox(slide, Inches(9.5), Inches(1.8), Inches(3.0), Inches(2.0))
    tf = tb_num.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = _make_run(p, "#1")
    run.font.size = Pt(72)
    run.font.color.rgb = CRITICAL_COLOR
    run.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    run2 = _make_run(p2, "OWASP Top 10\nfor LLM Apps")
    run2.font.size = Pt(14)
    run2.font.color.rgb = ACCENT_COLOR
    run2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    run3 = _make_run(p3, "(LLM01)")
    run3.font.size = Pt(12)
    run3.font.color.rgb = SUBTITLE_COLOR
    p3.alignment = PP_ALIGN.CENTER
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~45s] Define prompt injection clearly. Emphasize the SQL injection analogy \u2014 audience likely familiar with web security. Point to the #1 OWASP ranking to establish severity.")

    # --- Slide 4: Attack Taxonomy - Direct Injection (short rich bullets) ---
    print("Creating slide 4: Direct Injection")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Attack Taxonomy: Direct Injection")
    direct = ATTACK_TAXONOMY[0]
    _add_rich_bullets(slide, [
        ("Jailbreaking", "persona manipulation bypasses safety alignment (DAN, role-play scenarios)"),
        ("Goal Hijacking", "redirects model objective via competing instructions ('ignore previous...')"),
        ("Prompt Leaking", "extracts system prompts, revealing guardrails for targeted follow-up attacks"),
    ], font_size=18)
    # Subtitle context
    tb = _add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8))
    _set_text(tb, "Earliest and most widely studied class \u2014 exploits the inability to distinguish instructions from data",
              font_size=13, color=DIM_COLOR)
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~45s] Cover Direct Injection. Jailbreaking: DAN persona bypasses content policies. Goal hijacking: 'ignore previous instructions' redirects the model's task. Prompt leaking: extracting system prompts reveals application guardrails.")

    # --- Slide 5: Attack Taxonomy - Indirect Injection (short rich bullets) ---
    print("Creating slide 5: Indirect Injection")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Attack Taxonomy: Indirect Injection")
    _add_rich_bullets(slide, [
        ("RAG Poisoning", "malicious instructions in knowledge base documents execute on retrieval"),
        ("Web Content Injection", "hidden payloads in web pages target LLM browsing agents and summarizers"),
        ("Email & Document Injection", "instructions in emails/PDFs trigger unauthorized AI assistant actions"),
    ], font_size=18)
    # Subtitle context
    tb = _add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8))
    _set_text(tb, "Most dangerous vector for LLM-integrated apps \u2014 attacker never interacts with the model directly (Greshake et al., 2023)",
              font_size=13, color=DIM_COLOR)
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~45s] Cover Indirect Injection. Key insight: attacker plants payload in external content the model will retrieve. RAG poisoning is especially dangerous for enterprise knowledge bases. Email injection can trigger actions like forwarding data to attacker.")

    # --- Slide 6: Attack Taxonomy Figure (with annotation) ---
    print("Creating slide 6: Taxonomy figure")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Attack Taxonomy Visualization")
    # Annotation subtitle
    tb = _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4))
    _set_text(tb, "Six categories progressing from well-understood (direct) to emerging threats (autonomous propagation)",
              font_size=13, color=DIM_COLOR)
    _add_figure(slide, "attack_taxonomy.png", left=Inches(2), top=Inches(1.8),
                width=Inches(9), height=Inches(5.0))
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~30s] Walk through the taxonomy diagram. Point out the hierarchy from category to specific technique. Note how complexity increases left to right, reflecting the evolution of LLM capabilities.")

    # --- Slide 7: Critical Risks (with severity color) ---
    print("Creating slide 7: Critical Risks")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Real-World Risks: Critical Impact")
    critical = [r for r in RISK_CATEGORIES if r.get("severity") == "critical"]
    items = []
    for r in critical[:4]:
        examples = r.get("examples", [])
        ex_short = examples[0].split(" via ")[0] if examples else ""
        items.append((r['category'], ex_short))
    _add_rich_bullets(slide, items, font_size=17, label_color=CRITICAL_COLOR)
    # Severity indicator
    tb = _add_textbox(slide, Inches(0.8), Inches(5.8), Inches(4.0), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = _make_run(p, "\u25cf CRITICAL SEVERITY")
    run.font.size = Pt(12)
    run.font.color.rgb = CRITICAL_COLOR
    run.font.bold = True
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~45s] Emphasize data exfiltration and unauthorized actions as highest-severity outcomes. Give the markdown image exfiltration example \u2014 it's vivid and memorable: a poisoned doc encodes data into a URL rendered as an image tag.")

    # --- Slide 8: Broader Risks (with severity color) ---
    print("Creating slide 8: Broader Risks")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Risks: Broader Impact Categories")
    other = [r for r in RISK_CATEGORIES if r.get("severity") != "critical"]
    items = []
    for r in other[:5]:
        short_desc = r['description'].split('. ')[0]
        if len(short_desc) > 100:
            short_desc = short_desc[:97] + "..."
        items.append((r['category'], short_desc))
    _add_rich_bullets(slide, items, font_size=17, label_color=HIGH_COLOR)
    # Severity indicator
    tb = _add_textbox(slide, Inches(0.8), Inches(5.8), Inches(4.0), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = _make_run(p, "\u25cf HIGH SEVERITY")
    run.font.size = Pt(12)
    run.font.color.rgb = HIGH_COLOR
    run.font.bold = True
    # Key insight callout
    tb2 = _add_textbox(slide, Inches(5.0), Inches(5.8), Inches(7.5), Inches(0.5))
    _set_text(tb2, "Impact scales with agent permission scope \u2014 more tools = more risk",
              font_size=13, color=DIM_COLOR)
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~30s] Cover remaining risk categories briefly. Content manipulation is insidious because users trust AI output. Supply chain propagation means one compromised system can cascade to others. Highlight that impact scales with agent permission scope.")

    # --- Slide 9: Detection - Static & Learned (split 1/2) ---
    print("Creating slide 9: Detection (Static & Learned)")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Detection Techniques: Pattern-Based")
    _add_rich_bullets(slide, [
        ("Heuristic/Rule-Based", "keyword and regex matching against known injection patterns (Rebuff, WAF rules)"),
        ("ML Classification", "trained classifiers flag adversarial inputs via learned features (Lakera Guard, embedding similarity)"),
        ("Perplexity Analysis", "statistical anomaly detection flags unusual token distributions in adversarial suffixes"),
    ], font_size=17, top=Inches(1.5))
    # Tradeoff callout
    tb = _add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = _make_run(p, "Tradeoff: ")
    run.font.size = Pt(14)
    run.font.color.rgb = ACCENT_COLOR
    run.font.bold = True
    run2 = _make_run(p, "Fast and low-cost, but reactive \u2014 novel attacks consistently evade pattern-based detection")
    run2.font.size = Pt(14)
    run2.font.color.rgb = SUBTITLE_COLOR
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~30s] Cover the three pattern-based detection approaches. Heuristic: fast but trivially bypassed via paraphrasing. ML: better coverage but training-data dependent. Perplexity: good for machine-generated adversarial suffixes but not human-crafted prompts.")

    # --- Slide 10: Detection - Active & Model-Based (split 2/2) ---
    print("Creating slide 10: Detection (Active & Model-Based)")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Detection Techniques: Active Defense")
    _add_rich_bullets(slide, [
        ("Canary Tokens", "secret strings planted in system prompt \u2014 if they appear in output, injection succeeded"),
        ("LLM-as-Judge", "secondary model evaluates whether input attempts to subvert intended behavior"),
    ], font_size=17, top=Inches(1.5))
    # Tradeoff callout
    tb = _add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = _make_run(p, "Tradeoff: ")
    run.font.size = Pt(14)
    run.font.color.rgb = ACCENT_COLOR
    run.font.bold = True
    run2 = _make_run(p, "Smarter detection, but higher cost \u2014 LLM-as-judge adds latency and is itself vulnerable to injection")
    run2.font.size = Pt(14)
    run2.font.color.rgb = SUBTITLE_COLOR
    # Key takeaway
    p2 = tf.add_paragraph()
    p2.space_before = Pt(16)
    run3 = _make_run(p2, "\u2192 No single detection method is sufficient \u2014 defense in depth is required")
    run3.font.size = Pt(16)
    run3.font.color.rgb = HEADING_COLOR
    run3.font.bold = True
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~30s] Canary tokens: proactive but only catch context leakage. LLM-as-judge: most capable detector but doubles inference cost and the judge itself can be fooled. Key message: no single method works, you need to layer multiple approaches.")

    # --- Slide 11: Defense Flow Figure (with annotation) ---
    print("Creating slide 11: Defense flow figure")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Detection & Defense Pipeline")
    tb = _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4))
    _set_text(tb, "Multi-layered pipeline: input sanitization \u2192 detection \u2192 model processing \u2192 output filtering \u2192 monitoring",
              font_size=13, color=DIM_COLOR)
    _add_figure(slide, "defense_flow.png", left=Inches(2), top=Inches(1.8),
                width=Inches(9), height=Inches(5.0))
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~30s] Walk through the defense pipeline diagram. Each layer catches different attack types. Input sanitization stops known patterns. Detection flags anomalies. Output filtering catches successful attacks. Monitoring provides feedback loop.")

    # --- Slide 12: Prevention Strategies (rich bullets) ---
    print("Creating slide 12: Prevention")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Prevention Strategies")
    prev_items = [
        ("Instruction Hierarchy", "model-level privilege tiers \u2014 system > user > tool (Wallace et al., 2024)"),
        ("Input Sanitization", "strip injection patterns, normalize Unicode, enforce length limits"),
        ("Output Validation", "scan for data leakage, enforce format constraints, check canary tokens"),
        ("Sandboxing & Isolation", "least-privilege tool access, user approval gates for destructive actions"),
        ("Guardrail Frameworks", "programmable middleware enforcing conversation policies (NeMo, Lakera)"),
    ]
    _add_rich_bullets(slide, prev_items, font_size=16, top=Inches(1.4), height=Inches(5.0))
    # Takeaway
    tb = _add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = _make_run(p, "\u2192 No single strategy is complete \u2014 layer instruction hierarchy + sanitization + output filtering")
    run.font.size = Pt(14)
    run.font.color.rgb = ACCENT_COLOR
    run.font.bold = True
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~45s] Cover the layered defense approach. Instruction hierarchy is the most promising model-level defense but only model providers can implement it. App developers must layer sanitization + guardrails + output validation. System architects enforce sandboxing and least privilege. Emphasize that no single strategy is complete.")

    # --- Slide 13: Defense Architecture Figure (with annotation) ---
    print("Creating slide 13: Defense architecture figure")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Layered Defense Architecture")
    tb = _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4))
    _set_text(tb, "Model-level, application-level, and system-level controls working in coordination",
              font_size=13, color=DIM_COLOR)
    _add_figure(slide, "defense_architecture.png", left=Inches(2), top=Inches(1.8),
                width=Inches(9), height=Inches(5.0))
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~30s] Show the architecture diagram. Model providers handle instruction hierarchy. App developers deploy guardrails and filtering. System architects enforce sandboxing and least privilege. All three levels must coordinate for effective defense.")

    # --- Slide 14: Key References & Conclusion (fewer refs, larger font) ---
    print("Creating slide 14: References & Conclusion")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    _add_heading(slide, "Key References")
    # Show only 3 foundational references at readable size
    key_refs = [
        REFERENCES[0],  # Abdallah et al. (2025) - comprehensive review
        REFERENCES[7],  # Greshake et al. (2023) - indirect injection
        REFERENCES[15], # Wallace et al. (2024) - instruction hierarchy
    ]
    ref_items = [ref[:150] + "..." if len(ref) > 150 else ref for ref in key_refs]
    _add_bullets(slide, ref_items, font_size=14, top=Inches(1.3), height=Inches(3.0))
    # Source count note
    tb = _add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5))
    _set_text(tb, "Full reference list: 18 sources spanning academic papers, industry reports, and government standards (2023\u20132025)",
              font_size=12, color=DIM_COLOR)
    # Prominent conclusion line
    tb2 = _add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.2))
    tf = tb2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = _make_run(p, "Prompt injection remains an open problem \u2014\ndefense in depth is essential.")
    run.font.size = Pt(24)
    run.font.color.rgb = ACCENT_COLOR
    run.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[~30s] Highlight the three foundational references. Abdallah et al. for comprehensive taxonomy. Greshake et al. for indirect injection formalization. Wallace et al. for instruction hierarchy defense. Close with the key takeaway: no silver bullet, defense in depth required.")

    # --- Slide 15: Questions ---
    print("Creating slide 15: Questions")
    slide_num += 1
    slide = prs.slides.add_slide(blank_layout)
    _set_dark_bg(slide)
    # Large "Questions?" centered
    tb = _add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(2.0))
    _set_text(tb, "Questions?", font_size=48, color=HEADING_COLOR, bold=True,
              alignment=PP_ALIGN.CENTER)
    # Takeaway
    tb2 = _add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(1.0))
    _set_text(tb2, "Prompt injection is the #1 LLM vulnerability \u2014 no single defense is sufficient,\nbut layered strategies provide meaningful risk reduction.",
              font_size=18, color=SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
    # Contact / affiliation
    tb3 = _add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6))
    _set_text(tb3, f"{AUTHOR} \u2022 Christopher Newport University \u2022 COVACCI",
              font_size=14, color=DIM_COLOR, alignment=PP_ALIGN.CENTER)
    _add_slide_number(slide, slide_num)
    _add_notes(slide, "[End] Thank audience. Invite questions. If no questions, reiterate the defense-in-depth message and that prompt injection remains an active research area.")

    # Save
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_path = OUTPUT_DIR / "presentation.pptx"
    prs.save(str(output_path))
    print(f"\nPresentation saved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    generate_presentation()

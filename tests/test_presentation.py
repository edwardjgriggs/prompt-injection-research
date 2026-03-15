"""Tests for presentation generation output validation."""

import os
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

OUTPUT_PATH = Path(__file__).resolve().parent.parent / "output" / "presentation.pptx"

EXPECTED_FIGURES = ["attack_taxonomy.png", "defense_flow.png", "defense_architecture.png"]


@pytest.fixture(scope="module")
def prs():
    """Load the generated presentation once for all tests."""
    assert OUTPUT_PATH.exists(), f"Presentation not found at {OUTPUT_PATH}"
    return Presentation(str(OUTPUT_PATH))


class TestPresentationExists:
    def test_output_file_exists(self):
        assert OUTPUT_PATH.exists(), "output/presentation.pptx must exist"

    def test_file_is_valid_pptx(self, prs):
        # If we got here, python-pptx parsed it successfully
        assert prs is not None


class TestSlideCount:
    def test_slide_count_in_range(self, prs):
        count = len(prs.slides)
        assert 12 <= count <= 15, f"Expected 12-15 slides, got {count}"


class TestDarkBackground:
    def test_every_slide_has_dark_background(self, prs):
        for i, slide in enumerate(prs.slides):
            bg = slide.background
            fill = bg.fill
            # Check that background fill is solid and dark
            assert fill.type is not None, f"Slide {i+1} has no background fill"
            color = fill.fore_color.rgb
            r, g, b = color[0], color[1], color[2]
            # Dark means each channel <= 80
            assert r <= 80 and g <= 80 and b <= 80, (
                f"Slide {i+1} background RGB({r},{g},{b}) is not dark enough"
            )


class TestSpeakerNotes:
    def test_content_slides_have_speaker_notes(self, prs):
        """All slides except the first (title) should have speaker notes."""
        for i, slide in enumerate(prs.slides):
            if i == 0:
                continue  # skip title slide
            notes = slide.notes_slide.notes_text_frame.text.strip()
            assert len(notes) > 0, f"Slide {i+1} has no speaker notes"


class TestFigureEmbedding:
    def test_presentation_contains_figures(self, prs):
        """Check that expected figures are embedded as images somewhere in the deck."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        image_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
        assert image_count >= len(EXPECTED_FIGURES), (
            f"Expected at least {len(EXPECTED_FIGURES)} embedded images, found {image_count}"
        )
